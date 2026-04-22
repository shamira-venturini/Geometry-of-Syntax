import argparse
from pathlib import Path
from typing import Iterable, List, Sequence

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = REPO_ROOT / "docs"
DEFAULT_OUTPUT = DOCS_DIR / "experiment_1_methods_results_strict_control.pptx"

EXP1A_SUMMARY = (
    REPO_ROOT
    / "behavioral_results"
    / "experiment-1"
    / "experiment-1a"
    / "transitive_token_profiles"
    / "transitive_item_summary.csv"
)
EXP1A_PAIRED = (
    REPO_ROOT
    / "behavioral_results"
    / "experiment-1"
    / "experiment-1a"
    / "transitive_token_profiles"
    / "stats"
    / "paired_effects.csv"
)
EXP1B_ROOT = (
    REPO_ROOT
    / "behavioral_results"
    / "experiment-1"
    / "experiment-1b"
    / "processing_experiment_1b_gpt2large_v3_strict-control"
)
EXP1B_CORE_SUMMARY = EXP1B_ROOT / "processing_1b_core_core_lexically_controlled" / "summary.csv"
EXP1B_CORE_STATS = EXP1B_ROOT / "processing_1b_core_core_lexically_controlled" / "stats.csv"
EXP1B_JABBER_SUMMARY = EXP1B_ROOT / "processing_1b_jabberwocky_jabberwocky" / "summary.csv"
EXP1B_JABBER_STATS = EXP1B_ROOT / "processing_1b_jabberwocky_jabberwocky" / "stats.csv"
EXP1B_PRIMING = EXP1B_ROOT / "priming_framed_results.csv"
EXP1B_OVERLAP_ROOT = (
    REPO_ROOT
    / "behavioral_results"
    / "experiment-1"
    / "experiment-1b"
    / "processing_experiment_1b_gpt2large_v1_lexical-overlap"
)
EXP1B_OVERLAP_CORE_SUMMARY = EXP1B_OVERLAP_ROOT / "processing_1b_core_core" / "summary.csv"
EXP1B_OVERLAP_CORE_STATS = EXP1B_OVERLAP_ROOT / "processing_1b_core_core" / "stats.csv"
EXP1B_OVERLAP_PRIMING = EXP1B_OVERLAP_ROOT / "priming_framed_results.csv"
EXP1B_CORE_META = EXP1B_ROOT / "processing_1b_core_core_lexically_controlled" / "metadata.json"
EXP1B_JABBER_META = EXP1B_ROOT / "processing_1b_jabberwocky_jabberwocky" / "metadata.json"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(247, 244, 236)
NAVY = RGBColor(24, 41, 74)
TEAL = RGBColor(39, 118, 110)
ORANGE = RGBColor(191, 113, 63)
SAND = RGBColor(232, 221, 201)
TEXT = RGBColor(36, 38, 43)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(94, 96, 105)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Build a PPTX deck summarizing Experiment 1a and 1b methods/results."
    )
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    return parser.parse_args()


def fmt_float(value: float, digits: int = 4) -> str:
    return f"{value:.{digits}f}"


def fmt_effect(value: float) -> str:
    return f"{value:.3f}"


def fmt_p(value: float) -> str:
    if value == 0.0:
        return "<1e-300"
    if value < 0.001:
        return f"{value:.2e}"
    return f"{value:.3f}"


def set_cell_text(cell, text: str, *, font_size: int = 12, bold: bool = False, color: RGBColor = TEXT) -> None:
    cell.text = ""
    paragraph = cell.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    paragraph.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.color.rgb = NAVY

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(8.9), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(9.3), Inches(0.2), Inches(3.4), Inches(0.28))
        p = sub_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Aptos"
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(214, 222, 235)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(12.3), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED


def add_bullet_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    bullets: Sequence[str],
    accent: RGBColor,
) -> None:
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = accent
    frame.line.width = Pt(1.5)

    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(h))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.color.rgb = accent

    title_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.12), Inches(w - 0.3), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = accent

    body = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.55), Inches(w - 0.35), Inches(h - 0.68))
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.space_after = Pt(7)
        run = paragraph.add_run()
        run.text = bullet
        run.font.name = "Aptos"
        run.font.size = Pt(15)
        run.font.color.rgb = TEXT


def add_note_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    lines: Sequence[str],
    accent: RGBColor,
) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = SAND
    box.line.color.rgb = accent
    box.line.width = Pt(1.25)

    title_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.3), Inches(0.25))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = accent

    text_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.42), Inches(w - 0.28), Inches(h - 0.55))
    tf = text_box.text_frame
    first = True
    for line in lines:
        paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        paragraph.space_after = Pt(6)
        run = paragraph.add_run()
        run.text = line
        run.font.name = "Aptos"
        run.font.size = Pt(14)
        run.font.color.rgb = TEXT


def add_table(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    columns: Sequence[str],
    rows: Sequence[Sequence[str]],
    accent: RGBColor,
    font_size: int = 12,
) -> None:
    shape = slide.shapes.add_table(len(rows) + 1, len(columns), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table

    for idx, heading in enumerate(columns):
        cell = table.cell(0, idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent
        set_cell_text(cell, heading, font_size=font_size, bold=True, color=WHITE)

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 else BG
            set_cell_text(cell, value, font_size=font_size)


def exp1a_results() -> List[List[str]]:
    summary = pd.read_csv(EXP1A_SUMMARY)
    paired = pd.read_csv(EXP1A_PAIRED)

    rows: List[List[str]] = []
    for condition in ["CORE", "jabberwocky"]:
        active = summary[(summary["condition"] == condition) & (summary["target_structure"] == "active")].iloc[0]
        passive = summary[(summary["condition"] == condition) & (summary["target_structure"] == "passive")].iloc[0]
        diff = paired[(paired["condition"] == condition) & (paired["metric"] == "sentence_pe_mean")].iloc[0]
        rows.append(
            [
                "CORE" if condition == "CORE" else "Jabberwocky",
                fmt_float(active["mean_sentence_pe_mean"]),
                fmt_float(passive["mean_sentence_pe_mean"]),
                fmt_float(diff["mean_diff"]),
                f"[{fmt_float(diff['bootstrap_ci95_low'])}, {fmt_float(diff['bootstrap_ci95_high'])}]",
                fmt_effect(diff["effect_size_dz"]),
            ]
        )
    return rows


def exp1b_overall_results() -> List[List[str]]:
    rows: List[List[str]] = []
    configs = [
        ("Core strict-control", EXP1B_CORE_SUMMARY, EXP1B_CORE_STATS),
        ("Jabberwocky strict-control", EXP1B_JABBER_SUMMARY, EXP1B_JABBER_STATS),
    ]
    for label, summary_path, stats_path in configs:
        summary = pd.read_csv(summary_path)
        stats = pd.read_csv(stats_path)
        active = summary[summary["prime_condition"] == "active"].iloc[0]
        passive = summary[summary["prime_condition"] == "passive"].iloc[0]
        choice_stat = stats[
            (stats["condition_a"] == "active")
            & (stats["condition_b"] == "passive")
            & (stats["metric"] == "passive_choice_delta")
        ].iloc[0]
        logprob_stat = stats[
            (stats["condition_a"] == "active")
            & (stats["condition_b"] == "passive")
            & (stats["metric"] == "logprob_delta")
        ].iloc[0]
        rows.append(
            [
                label,
                fmt_float(active["passive_choice_rate"]),
                fmt_float(passive["passive_choice_rate"]),
                fmt_float(choice_stat["mean_diff_b_minus_a"]),
                fmt_p(choice_stat["t_p_two_sided"]),
                fmt_float(logprob_stat["mean_diff_b_minus_a"]),
            ]
        )
    return rows


def exp1b_baseline_rows() -> List[List[str]]:
    priming = pd.read_csv(EXP1B_PRIMING)
    ordered = []
    for condition in ["core_core", "jabberwocky_jabberwocky"]:
        for baseline in ["no_prime_eos", "no_prime_empty", "filler"]:
            row = priming[(priming["condition"] == condition) & (priming["baseline"] == baseline)].iloc[0]
            ordered.append(
                [
                    "Core" if condition == "core_core" else "Jabberwocky",
                    baseline,
                    fmt_float(row["baseline_passive_choice_rate"]),
                    fmt_float(row["active_choice_priming"]),
                    fmt_float(row["passive_choice_priming"]),
                    fmt_float(row["imbalance_choice_passive_minus_active"]),
                ]
            )
    return ordered


def exp1b_overlap_overall_rows() -> List[List[str]]:
    rows: List[List[str]] = []
    configs = [
        ("Core strict-control", EXP1B_CORE_SUMMARY, EXP1B_CORE_STATS),
        ("Core lexical-overlap", EXP1B_OVERLAP_CORE_SUMMARY, EXP1B_OVERLAP_CORE_STATS),
    ]
    for label, summary_path, stats_path in configs:
        summary = pd.read_csv(summary_path)
        stats = pd.read_csv(stats_path)
        active = summary[summary["prime_condition"] == "active"].iloc[0]
        passive = summary[summary["prime_condition"] == "passive"].iloc[0]
        choice_stat = stats[
            (stats["condition_a"] == "active")
            & (stats["condition_b"] == "passive")
            & (stats["metric"] == "passive_choice_delta")
        ].iloc[0]
        logprob_stat = stats[
            (stats["condition_a"] == "active")
            & (stats["condition_b"] == "passive")
            & (stats["metric"] == "logprob_delta")
        ].iloc[0]
        rows.append(
            [
                label,
                fmt_float(active["passive_choice_rate"]),
                fmt_float(passive["passive_choice_rate"]),
                fmt_float(choice_stat["mean_diff_b_minus_a"]),
                fmt_float(logprob_stat["mean_diff_b_minus_a"]),
            ]
        )
    return rows


def exp1b_overlap_filler_rows() -> List[List[str]]:
    strict = pd.read_csv(EXP1B_PRIMING)
    overlap = pd.read_csv(EXP1B_OVERLAP_PRIMING)
    strict_row = strict[(strict["condition"] == "core_core") & (strict["baseline"] == "filler")].iloc[0]
    overlap_row = overlap[(overlap["condition"] == "core_core") & (overlap["baseline"] == "filler")].iloc[0]
    return [
        [
            "Active priming",
            fmt_float(strict_row["active_choice_priming"]),
            fmt_float(overlap_row["active_choice_priming"]),
            fmt_float(overlap_row["active_choice_priming"] - strict_row["active_choice_priming"]),
        ],
        [
            "Passive priming",
            fmt_float(strict_row["passive_choice_priming"]),
            fmt_float(overlap_row["passive_choice_priming"]),
            fmt_float(overlap_row["passive_choice_priming"] - strict_row["passive_choice_priming"]),
        ],
        [
            "Imbalance (P-A)",
            fmt_float(strict_row["imbalance_choice_passive_minus_active"]),
            fmt_float(overlap_row["imbalance_choice_passive_minus_active"]),
            fmt_float(overlap_row["imbalance_choice_passive_minus_active"] - strict_row["imbalance_choice_passive_minus_active"]),
        ],
    ]


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.25))
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.color.rgb = NAVY

    title = slide.shapes.add_textbox(Inches(0.55), Inches(1.55), Inches(7.8), Inches(1.15))
    p = title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Experiment 1a and 1b\nMethods and Results"
    run.font.name = "Aptos Display"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = NAVY

    sub = slide.shapes.add_textbox(Inches(0.58), Inches(2.9), Inches(7.8), Inches(0.55))
    p = sub.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Structural priming in GPT-2 Large, updated with strict-control Experiment 1b outputs"
    run.font.name = "Aptos"
    run.font.size = Pt(18)
    run.font.color.rgb = MUTED

    add_note_box(
        slide,
        8.35,
        1.55,
        4.3,
        3.0,
        "Deck scope",
        [
            "Methods for Experiment 1a",
            "Results for Experiment 1a",
            "Methods for Experiment 1b",
            "Results for Experiment 1b",
            "1b results use the repaired v3 strict-control runs.",
        ],
        TEAL,
    )
    add_footer(slide, "Geometry-of-Syntax | Generated automatically on 2026-04-10")

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Experiment 1a: Methods", "Processing replication")
    add_bullet_box(
        slide,
        0.55,
        1.25,
        7.35,
        5.15,
        "Design",
        [
            "Sinclair-style teacher-forced structural priming in gpt2-large.",
            "Conditions: CORE prime -> CORE target and Jabberwocky prime -> Jabberwocky target.",
            "Materials: 15,000 transitive items per condition, with active/passive prime and target realizations.",
            "For each item, the target sentence is scored after a congruent versus incongruent prime.",
            "Primary metric: sentence_pe_mean (mean token priming effect across the target sentence).",
            "Inference: paired passive-minus-active contrasts with t-tests, sign-flip permutation tests, bootstrap CIs, and mixed models.",
        ],
        TEAL,
    )
    add_note_box(
        slide,
        8.15,
        1.25,
        4.55,
        2.4,
        "Run summary",
        [
            "Model: gpt2-large",
            "n = 15,000 per condition",
            "Key scripts: 2_transitive_token_priming.py, 3_summarize_transitive_priming.py, 5_analyze_transitive_statistics.py",
        ],
        ORANGE,
    )
    add_note_box(
        slide,
        8.15,
        3.95,
        4.55,
        2.15,
        "Interpretive role",
        [
            "Experiment 1a is the replication anchor for the project.",
            "The Jabberwocky condition is a matched extension, not part of Sinclair et al. 2022 itself.",
        ],
        ORANGE,
    )
    add_footer(slide, "Experiment 1a uses the canonical transitive_token_profiles outputs.")

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Experiment 1a: Results", "Primary metric = sentence_pe_mean")
    add_table(
        slide,
        0.6,
        1.25,
        8.25,
        2.25,
        ["Condition", "Active PE", "Passive PE", "Passive-Active", "95% CI", "d_z"],
        exp1a_results(),
        TEAL,
        font_size=12,
    )
    add_note_box(
        slide,
        9.05,
        1.28,
        3.9,
        2.0,
        "Headline",
        [
            "CORE reproduces the directional Sinclair-style pattern: both structures show positive priming, with stronger passive priming.",
            "Jabberwocky shows a much larger passive-skewed asymmetry.",
        ],
        ORANGE,
    )
    add_bullet_box(
        slide,
        0.6,
        3.9,
        12.35,
        2.35,
        "Takeaways",
        [
            "CORE: active PE = 0.0312, passive PE = 0.1746, passive-active = 0.1434, 95% CI [0.1375, 0.1493], d_z = 0.389.",
            "Jabberwocky: active PE = -0.1121, passive PE = 0.3854, passive-active = 0.4975, 95% CI [0.4921, 0.5029], d_z = 1.482.",
            "The passive-over-active effect is therefore present in both conditions and much larger in Jabberwocky.",
        ],
        ORANGE,
    )
    add_footer(slide, "Source: experiment-1a/transitive_item_summary.csv and stats/paired_effects.csv")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Experiment 1b: Methods", "Controlled processing follow-up")
    add_bullet_box(
        slide,
        0.55,
        1.25,
        7.4,
        5.25,
        "Design",
        [
            "Teacher-forced active/passive target competition for the same event, scored under multiple prime conditions.",
            "Prime conditions: active, passive, no_prime_eos, no_prime_empty, and filler.",
            "For each item, both the full active and full passive target sentences are scored.",
            "Primary outputs: passive-choice rate and passive-minus-active mean target log probability.",
            "The analysis emphasizes baseline decomposition, not only the active-vs-passive prime contrast.",
        ],
        ORANGE,
    )
    add_note_box(
        slide,
        8.15,
        1.25,
        4.55,
        2.55,
        "Strict-control v3",
        [
            "Core corpus repaired to remove prime-target lexical overlap.",
            "Core lexical audit: 0 shared verbs, 0 shared nouns.",
            "Jabberwocky rerun matched at n = 2080 with seed 13.",
        ],
        TEAL,
    )
    add_note_box(
        slide,
        8.15,
        4.05,
        4.55,
        2.1,
        "Conditions used in deck",
        [
            "Core: repaired lexically controlled run",
            "Jabberwocky: completed strict-control rerun",
            "Model: gpt2-large",
        ],
        TEAL,
    )
    add_footer(slide, "Results slides below use processing_experiment_1b_gpt2large_v3_strict-control.")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Experiment 1b: Methods", "Corpus revisions across batches")
    add_bullet_box(
        slide,
        0.55,
        1.18,
        7.55,
        5.25,
        "What changed in the 1b materials",
        [
            "We moved from the original constrained CORE transitive set to a reduced reversible counterbalanced corpus where eligible nouns cycle equally often through AGENT and PATIENT roles within each verb block.",
            "After finding a lexical-overlap bug, we built a repaired CORE strict-control corpus in which prime and target share neither the active-form verb nor either noun.",
            "For the repaired strict-control comparison, both CORE and Jabberwocky conditions were matched at n = 2080 items rather than leaving Jabberwocky at the full 15,000-row source corpus.",
            "The filler baseline was tightened from a generic shared sentence context to domain-matched filler pools: real-word fillers for CORE and nonce fillers for Jabberwocky.",
            "The old lexical-overlap CORE corpus was retained only as a comparison condition, while the presentation deck uses the repaired lexically controlled version as the main evidence.",
        ],
        TEAL,
    )
    add_note_box(
        slide,
        8.35,
        1.2,
        4.2,
        2.65,
        "Strict-control audit",
        [
            "CORE repaired corpus audit: 0 shared prime-target verbs",
            "CORE repaired corpus audit: 0 shared prime-target nouns",
            "Jabberwocky strict-control rerun: 2080 sampled items, paired_indices alignment, Jabberwocky filler pool",
        ],
        ORANGE,
    )
    add_note_box(
        slide,
        8.35,
        4.15,
        4.2,
        2.0,
        "Why it matters",
        [
            "These corpus fixes reduce lexical boost confounds and make the baseline decomposition interpretable.",
            "They also explain why the strict-control CORE effect is much smaller than the older overlap-heavy runs.",
        ],
        ORANGE,
    )
    add_footer(
        slide,
        "Sources: project_status_handoff.md, 11_generate_counterbalanced_constrained_transitive.py, 17_build_lexically_controlled_core_transitive.py",
    )

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Experiment 1b: Results", "Overall active-vs-passive prime contrast")
    add_table(
        slide,
        0.6,
        1.25,
        8.55,
        2.25,
        ["Condition", "Passive Choice (A)", "Passive Choice (P)", "Choice Shift", "Choice p", "Logprob Shift"],
        exp1b_overall_results(),
        ORANGE,
        font_size=12,
    )
    add_note_box(
        slide,
        9.35,
        1.25,
        3.55,
        2.15,
        "Key point",
        [
            "Both strict-control conditions still show reliable priming in the active-vs-passive contrast.",
            "But both also begin from strongly passive-biased baselines.",
        ],
        TEAL,
    )
    add_bullet_box(
        slide,
        0.6,
        3.95,
        12.3,
        2.3,
        "Readout",
        [
            "Core strict-control: passive choice rises from 0.9779 after active primes to 1.0000 after passive primes (shift = +0.0221, p = 9.24e-12; logprob shift = +0.6372).",
            "Jabberwocky strict-control: passive choice rises from 0.9543 to 0.9995 (shift = +0.0452, p = 1.08e-22; logprob shift = +0.2788).",
            "The lexical-overlap repair sharply reduced the core-core effect relative to older overlap-heavy runs.",
        ],
        TEAL,
    )
    add_footer(slide, "Choice p-values from active vs passive paired passive_choice_delta tests.")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Experiment 1b: Results", "Baseline decomposition of priming")
    add_table(
        slide,
        0.55,
        1.18,
        8.85,
        3.25,
        ["Condition", "Baseline", "Baseline Passive", "Active Priming", "Passive Priming", "Imbalance"],
        exp1b_baseline_rows(),
        TEAL,
        font_size=11,
    )
    add_note_box(
        slide,
        9.6,
        1.2,
        3.05,
        3.0,
        "How to read",
        [
            "Active priming = movement toward active relative to baseline.",
            "Passive priming = movement toward passive relative to baseline.",
            "Imbalance = passive priming - active priming.",
            "Positive imbalance means passive priming is stronger.",
        ],
        ORANGE,
    )
    add_bullet_box(
        slide,
        0.55,
        4.7,
        12.1,
        1.55,
        "Interpretation",
        [
            "Baseline passive-choice rates are very high in both conditions, so overall prime contrasts need baseline context.",
            "Core strict-control is small and baseline-sensitive: filler and no_prime_empty favor larger active than passive priming, but no_prime_eos slightly favors passive priming.",
            "Jabberwocky strict-control remains more consistently passive-skewed, especially against the no-prime baselines.",
        ],
        ORANGE,
    )
    add_footer(slide, "Imbalance values are based on passive-choice decomposition in priming_framed_results.csv")

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header(slide, "Experiment 1b: Results", "Lexical-overlap inflated the Core effect")
    add_table(
        slide,
        0.55,
        1.2,
        6.25,
        2.2,
        ["Core condition", "Passive Choice (A)", "Passive Choice (P)", "Choice Shift", "Logprob Shift"],
        exp1b_overlap_overall_rows(),
        ORANGE,
        font_size=11,
    )
    add_table(
        slide,
        0.55,
        3.9,
        6.25,
        1.9,
        ["Filler-baseline measure", "Strict", "Overlap", "Overlap - Strict"],
        exp1b_overlap_filler_rows(),
        TEAL,
        font_size=11,
    )
    add_note_box(
        slide,
        7.15,
        1.2,
        5.45,
        4.55,
        "Interpretation",
        [
            "Lexical overlap increased the Core active-vs-passive prime contrast from 0.0221 to 0.1654 on the binary choice metric.",
            "The corresponding logprob shift increased from 0.6372 to 1.4856.",
            "Against the filler baseline, active priming rose from 0.0159 to 0.1591, while passive priming stayed at 0.0063 on the binary choice measure.",
            "So the overlap bug mostly made active primes look much more effective against the same passive-skewed baseline, rather than raising passive priming or changing the baseline itself.",
        ],
        NAVY,
    )
    add_footer(
        slide,
        "Comparison uses v1 lexical-overlap Core vs v3 strict-control Core; lexical boost was a Core corpus issue.",
    )

    return prs


def main() -> None:
    args = parse_args()
    prs = build_presentation()
    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    print(f"Saved presentation to {args.output}")


if __name__ == "__main__":
    main()
